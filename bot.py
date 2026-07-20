#!/usr/bin/env python3
"""
PDF Converter Telegram Bot
Fayllarni (rasmlar, albomlar, docx, txt, xlsx, pptx, pdf) PDF ga va aksincha aylantiruvchi universal bot.
"""

import os
import sys
import io
import re
import glob
import asyncio
import logging
import tempfile
import chardet
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional
from dotenv import load_dotenv

from PIL import Image
import img2pdf
import fitz  # PyMuPDF
from pypdf import PdfReader, PdfWriter
import docx
import openpyxl
from pptx import Presentation
from pdf2docx import Converter

import reportlab
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
    PageBreak,
    KeepTogether,
    Image as RLImage,
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas

from telegram import (
    Update,
    InlineKeyboardButton,
    InlineKeyboardMarkup,
    InputMediaDocument,
)
from telegram.ext import (
    Application,
    CommandHandler,
    MessageHandler,
    CallbackQueryHandler,
    ConversationHandler,
    ContextTypes,
    filters,
)

# Logging sozlamalari
logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    level=logging.INFO,
)
logger = logging.getLogger(__name__)

load_dotenv()
TOKEN = os.getenv("BOT_TOKEN", "YOUR_BOT_TOKEN_HERE")

# States for ConversationHandler if needed
WAITING_PASSWORD = 1
WAITING_PAGES = 2
WAITING_WATERMARK = 3
WAITING_STAMP_IMG = 4

# Shriftlarni ro'yxatdan o'tkazish
FONT_NAME = "Helvetica"
FONT_BOLD = "Helvetica-Bold"

def init_fonts():
    global FONT_NAME, FONT_BOLD
    font_candidates = [
        ("C:/Windows/Fonts/arial.ttf", "C:/Windows/Fonts/arialbd.ttf"),
        ("C:/Windows/Fonts/calibri.ttf", "C:/Windows/Fonts/calibrib.ttf"),
        ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
        ("/usr/share/fonts/TTF/DejaVuSans.ttf", "/usr/share/fonts/TTF/DejaVuSans-Bold.ttf"),
        ("/Library/Fonts/Arial.ttf", "/Library/Fonts/Arial Bold.ttf"),
    ]
    
    for font_path, bold_path in font_candidates:
        if os.path.exists(font_path):
            try:
                pdfmetrics.registerFont(TTFont("CustomUTF8", font_path))
                FONT_NAME = "CustomUTF8"
                if os.path.exists(bold_path):
                    pdfmetrics.registerFont(TTFont("CustomUTF8-Bold", bold_path))
                    FONT_BOLD = "CustomUTF8-Bold"
                else:
                    FONT_BOLD = "CustomUTF8"
                logger.info(f"UTF-8 Shrift yuklandi: {font_path}")
                return
            except Exception as e:
                logger.warning(f"Shrift xatosi ({font_path}): {e}")

init_fonts()

# Import OpenCV for Smart Auto-Cropping & Enhancement
import cv2
import numpy as np

def order_points(pts: np.ndarray) -> np.ndarray:
    rect = np.zeros((4, 2), dtype="float32")
    s = pts.sum(axis=1)
    rect[0] = pts[np.argmin(s)]
    rect[2] = pts[np.argmax(s)]
    diff = np.diff(pts, axis=1)
    rect[1] = pts[np.argmin(diff)]
    rect[3] = pts[np.argmax(diff)]
    return rect

def four_point_transform(image: np.ndarray, pts: np.ndarray) -> np.ndarray:
    rect = order_points(pts)
    (tl, tr, br, bl) = rect

    widthA = np.sqrt(((br[0] - bl[0]) ** 2) + ((br[1] - bl[1]) ** 2))
    widthB = np.sqrt(((tr[0] - tl[0]) ** 2) + ((tr[1] - tl[1]) ** 2))
    maxWidth = max(int(widthA), int(widthB))

    heightA = np.sqrt(((tr[0] - br[0]) ** 2) + ((tr[1] - br[1]) ** 2))
    heightB = np.sqrt(((tl[0] - bl[0]) ** 2) + ((tl[1] - bl[1]) ** 2))
    maxHeight = max(int(heightA), int(heightB))

    dst = np.array([
        [0, 0],
        [maxWidth - 1, 0],
        [maxWidth - 1, maxHeight - 1],
        [0, maxHeight - 1]], dtype="float32")

    M = cv2.getPerspectiveTransform(rect, dst)
    warped = cv2.warpPerspective(image, M, (maxWidth, maxHeight))
    return warped

def auto_crop_document_sync(image_path: str, output_path: str) -> str:
    """OpenCV yordamida hujjat burchaklarini topib, stoldan ajratish va tekislash (Perspective Warp)"""
    image = cv2.imread(image_path)
    if image is None:
        return image_path

    orig = image.copy()
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    blurred = cv2.GaussianBlur(gray, (5, 5), 0)
    edged = cv2.Canny(blurred, 75, 200)

    contours, _ = cv2.findContours(edged.copy(), cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)
    contours = sorted(contours, key=cv2.contourArea, reverse=True)[:5]

    doc_cnt = None
    for c in contours:
        peri = cv2.arcLength(c, True)
        approx = cv2.approxPolyDP(c, 0.02 * peri, True)
        if len(approx) == 4:
            doc_cnt = approx
            break

    if doc_cnt is not None:
        warped = four_point_transform(orig, doc_cnt.reshape(4, 2))
    else:
        warped = orig

    cv2.imwrite(output_path, warped)
    return output_path

def enhance_document_scan_sync(image_path: str, output_path: str, mode: str = "magic_color") -> str:
    """Hujjat rasmini yorqin skaner holatiga keltirish"""
    image = cv2.imread(image_path)
    if image is None:
        return image_path

    if mode == "magic_color":
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        bg = cv2.medianBlur(gray, 21)
        diff = 255 - cv2.absdiff(gray, bg)
        norm_diff = cv2.normalize(diff, None, alpha=0, beta=255, norm_type=cv2.NORM_MINMAX)
        result = cv2.cvtColor(norm_diff, cv2.COLOR_GRAY2BGR)
    elif mode == "bw":
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        result = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2)
    else:
        result = image

    cv2.imwrite(output_path, result)
    return output_path
TEXT_EXTENSIONS = {".txt", ".md", ".log", ".json", ".csv", ".py", ".js", ".html", ".css", ".xml"}
DOCX_EXTENSIONS = {".docx"}
EXCEL_EXTENSIONS = {".xlsx", ".xls"}
PPTX_EXTENSIONS = {".pptx", ".ppt"}

def is_image_file(filename: str) -> bool:
    return Path(filename).suffix.lower() in IMAGE_EXTENSIONS

def is_text_file(filename: str) -> bool:
    return Path(filename).suffix.lower() in TEXT_EXTENSIONS

def is_docx_file(filename: str) -> bool:
    return Path(filename).suffix.lower() in DOCX_EXTENSIONS

def is_excel_file(filename: str) -> bool:
    return Path(filename).suffix.lower() in EXCEL_EXTENSIONS

def is_pptx_file(filename: str) -> bool:
    return Path(filename).suffix.lower() in PPTX_EXTENSIONS

def render_progress_bar(percent: int) -> str:
    filled = int(percent / 10)
    bar = "█" * filled + "░" * (10 - filled)
    return f"[{bar}] {percent}%"

# Numbered Canvas class for Page X of Y footers
class NumberedCanvas(canvas.Canvas):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_page_number(num_pages)
            super().showPage()
        super().save()

    def draw_page_number(self, page_count):
        self.saveState()
        self.setFont(FONT_NAME, 9)
        self.setFillColor(colors.HexColor("#666666"))
        page_text = f"Sahifa {self._pageNumber} / {page_count}"
        self.drawRightString(A4[0] - 36, 25, page_text)
        self.drawString(36, 25, "PDF Converter Bot tomonidan yaratildi")
        self.setStrokeColor(colors.HexColor("#CCCCCC"))
        self.setLineWidth(0.5)
        self.line(36, 38, A4[0] - 36, 38)
        self.restoreState()


# ==========================================
# CONVERTER UTILITIES (CPU-bound Sync)
# ==========================================

def images_to_pdf_sync(image_paths: List[str], output_pdf_path: str) -> str:
    """Rasmlar ro'yxatini yagona PDF ga o'girish"""
    valid_paths = []
    temp_files_to_clean = []
    
    for path in image_paths:
        if not os.path.exists(path):
            continue
        try:
            with Image.open(path) as img:
                if img.mode in ("RGBA", "P", "LA", "CMYK"):
                    rgb_img = Image.new("RGB", img.size, (255, 255, 255))
                    if img.mode == "RGBA":
                        rgb_img.paste(img, mask=img.split()[3])
                    else:
                        rgb_img.paste(img.convert("RGB"))
                    
                    tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
                    rgb_img.save(tmp_file.name, "JPEG", quality=95)
                    tmp_file.close()
                    valid_paths.append(tmp_file.name)
                    temp_files_to_clean.append(tmp_file.name)
                else:
                    valid_paths.append(path)
        except Exception as e:
            logger.error(f"Rasm xatosi ({path}): {e}")

    if not valid_paths:
        raise ValueError("Yaroqli rasmlar topilmadi.")

    try:
        pdf_bytes = img2pdf.convert(valid_paths)
        with open(output_pdf_path, "wb") as f:
            f.write(pdf_bytes)
    finally:
        for tmp in temp_files_to_clean:
            if os.path.exists(tmp):
                try:
                    os.remove(tmp)
                except Exception:
                    pass

    return output_pdf_path


def images_to_pdf_a4_sync(image_paths: List[str], output_pdf_path: str) -> str:
    """Rasmlarni A4 sahifasiga moslab (A4 Fit proportsional scaling) PDF ga o'girish"""
    c = canvas.Canvas(output_pdf_path, pagesize=A4)
    page_w, page_h = A4[0], A4[1]
    processed_count = 0
    
    for path in image_paths:
        if not os.path.exists(path):
            continue
        try:
            with Image.open(path) as img:
                w, h = img.size
                scale = min(page_w / w, page_h / h)
                new_w, new_h = w * scale, h * scale
                x = (page_w - new_w) / 2
                y = (page_h - new_h) / 2
                
                c.drawImage(path, x, y, width=new_w, height=new_h)
                c.showPage()
                processed_count += 1
        except Exception as e:
            logger.error(f"Image to A4 error ({path}): {e}")

    if processed_count == 0:
        raise ValueError("Yaroqli rasmlar topilmadi.")

    c.save()
    return output_pdf_path


def text_to_pdf_sync(text_content: str, output_pdf_path: str, title: str = "Hujjat") -> str:
    """Matnni PDF ga o'girish"""
    doc = SimpleDocTemplate(
        output_pdf_path,
        pagesize=A4,
        leftMargin=36, rightMargin=36, topMargin=45, bottomMargin=50,
    )
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "DocTitle", parent=styles["Heading1"], fontName=FONT_BOLD, fontSize=18, leading=22, textColor=colors.HexColor("#1A365D"), spaceAfter=15,
    )
    body_style = ParagraphStyle(
        "DocBody", parent=styles["Normal"], fontName=FONT_NAME, fontSize=11, leading=15, textColor=colors.HexColor("#2D3748"), spaceAfter=8,
    )

    story = [Paragraph(title, title_style), Spacer(1, 10)]
    for line in text_content.split("\n"):
        clean_line = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").strip()
        if clean_line:
            story.append(Paragraph(clean_line, body_style))
        else:
            story.append(Spacer(1, 6))

    doc.build(story, canvasmaker=NumberedCanvas)
    return output_pdf_path


def docx_to_pdf_sync(docx_path: str, output_pdf_path: str) -> str:
    """.docx faylini PDF ga o'girish"""
    document = docx.Document(docx_path)
    doc = SimpleDocTemplate(
        output_pdf_path, pagesize=A4, leftMargin=36, rightMargin=36, topMargin=45, bottomMargin=50,
    )
    styles = getSampleStyleSheet()
    heading_style = ParagraphStyle(
        "DocxHeading", parent=styles["Heading1"], fontName=FONT_BOLD, fontSize=16, leading=20, textColor=colors.HexColor("#1A202C"), spaceBefore=12, spaceAfter=6,
    )
    body_style = ParagraphStyle(
        "DocxBody", parent=styles["Normal"], fontName=FONT_NAME, fontSize=10.5, leading=14.5, textColor=colors.HexColor("#2D3748"), spaceAfter=6,
    )

    story = []
    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            story.append(Spacer(1, 4))
            continue
        safe_text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        if para.style.name.startswith("Heading"):
            story.append(Paragraph(safe_text, heading_style))
        else:
            story.append(Paragraph(safe_text, body_style))

    for table in document.tables:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                row_data.append(Paragraph(cell_text, body_style))
            table_data.append(row_data)
        
        if table_data:
            t = Table(table_data)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#EDF2F7')),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#CBD5E0')),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                ('TOPPADDING', (0, 0), (-1, -1), 6),
            ]))
            story.append(Spacer(1, 8))
            story.append(t)
            story.append(Spacer(1, 8))

    doc.build(story, canvasmaker=NumberedCanvas)
    return output_pdf_path


def excel_to_pdf_sync(excel_path: str, output_pdf_path: str) -> str:
    """Excel (.xlsx) ni PDF ga o'girish"""
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    doc = SimpleDocTemplate(output_pdf_path, pagesize=A4, leftMargin=20, rightMargin=20, topMargin=40, bottomMargin=40)
    story = []
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("ExTitle", parent=styles["Heading1"], fontName=FONT_BOLD, fontSize=15, textColor=colors.HexColor("#2B6CB0"), spaceAfter=8)
    cell_style = ParagraphStyle("ExCell", parent=styles["Normal"], fontName=FONT_NAME, fontSize=8.5, leading=11)

    for sheet in wb.worksheets:
        story.append(Paragraph(f"📊 Varaq: {sheet.title}", title_style))
        story.append(Spacer(1, 6))
        table_data = []
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                row_cells = [Paragraph(str(c) if c is not None else "", cell_style) for c in row]
                table_data.append(row_cells)
        if table_data:
            t = Table(table_data)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#E2E8F0')),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#CBD5E0')),
                ('TOPPADDING', (0, 0), (-1, -1), 4),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
            ]))
            story.append(t)
            story.append(Spacer(1, 15))

    doc.build(story, canvasmaker=NumberedCanvas)
    return output_pdf_path


def pptx_to_pdf_sync(pptx_path: str, output_pdf_path: str) -> str:
    """PowerPoint (.pptx) ni PDF ga o'girish"""
    prs = Presentation(pptx_path)
    doc = SimpleDocTemplate(output_pdf_path, pagesize=A4, leftMargin=36, rightMargin=36, topMargin=40, bottomMargin=45)
    story = []
    styles = getSampleStyleSheet()
    h_style = ParagraphStyle("SlideH", parent=styles["Heading1"], fontName=FONT_BOLD, fontSize=16, textColor=colors.HexColor("#2C5282"), leading=20, spaceAfter=8)
    b_style = ParagraphStyle("SlideB", parent=styles["Normal"], fontName=FONT_NAME, fontSize=11, leading=15, spaceAfter=6)

    for idx, slide in enumerate(prs.slides):
        story.append(Paragraph(f"🎨 Slayd {idx + 1}", h_style))
        story.append(Spacer(1, 4))
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        safe_text = p.text.strip().replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                        story.append(Paragraph(safe_text, b_style))
        story.append(Spacer(1, 10))
        if idx < len(prs.slides) - 1:
            story.append(PageBreak())

    doc.build(story, canvasmaker=NumberedCanvas)
    return output_pdf_path


def pdf_to_docx_sync(input_pdf_path: str, output_docx_path: str) -> str:
    """PDF ni qayta Word (.docx) fayliga o'girish"""
    cv = Converter(input_pdf_path)
    cv.convert(output_docx_path)
    cv.close()
    return output_docx_path


def pdf_encrypt_sync(input_pdf_path: str, output_pdf_path: str, password: str) -> str:
    """PDF ga parol o'rnatish"""
    reader = PdfReader(input_pdf_path)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(password)
    with open(output_pdf_path, "wb") as f:
        writer.write(f)
    return output_pdf_path


def pdf_decrypt_sync(input_pdf_path: str, output_pdf_path: str, password: str) -> str:
    """PDF parolini olib tashlash"""
    reader = PdfReader(input_pdf_path)
    if reader.is_encrypted:
        success = reader.decrypt(password)
        if not success:
            raise ValueError("Noto'g'ri parol kiritildi!")
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    with open(output_pdf_path, "wb") as f:
        writer.write(f)
    return output_pdf_path


def pdf_delete_pages_sync(input_pdf_path: str, output_pdf_path: str, pages_to_remove: List[int]) -> str:
    """PDF dan ko'rsatilgan sahifalarni olib tashlash"""
    reader = PdfReader(input_pdf_path)
    writer = PdfWriter()
    remove_set = set(p - 1 for p in pages_to_remove)
    
    added_count = 0
    for idx, page in enumerate(reader.pages):
        if idx not in remove_set:
            writer.add_page(page)
            added_count += 1
            
    if added_count == 0:
        raise ValueError("Barcha sahifalar o'chirib tashlandi!")

    with open(output_pdf_path, "wb") as f:
        writer.write(f)
    return output_pdf_path


def pdf_add_watermark_sync(input_pdf_path: str, output_pdf_path: str, text: str) -> str:
    """PDF sahifalariga suv belgisi (Watermark) qo'shish"""
    doc = fitz.open(input_pdf_path)
    try:
        for page in doc:
            rect = page.rect
            page.insert_text(
                fitz.Point(rect.width / 4, rect.height / 2),
                text,
                fontsize=36,
                color=(0.7, 0.7, 0.7),
                rotate=0,
                overlay=True
            )
        doc.save(output_pdf_path)
    finally:
        doc.close()
    return output_pdf_path


# Import Translation and OCR
from deep_translator import GoogleTranslator

try:
    import pytesseract
    def init_tesseract():
        tesseract_paths = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            r"C:\Users\Leopard\AppData\Local\Programs\Tesseract-OCR\tesseract.exe",
            "/usr/bin/tesseract",
            "/usr/local/bin/tesseract",
        ]
        for p in tesseract_paths:
            if os.path.exists(p):
                pytesseract.pytesseract.tesseract_cmd = p
                logger.info(f"Tesseract OCR topildi: {p}")
                return
    init_tesseract()
except Exception as e:
    logger.warning(f"PyTesseract yuklanishda xatolik: {e}")
    pytesseract = None

def translate_text_sync(text: str, target_lang: str = "uz") -> str:
    """Matnni ko'rsatilgan tilga tarjima qilish"""
    if not text.strip():
        return text
    translator = GoogleTranslator(source="auto", target=target_lang)
    chunks = text.split("\n")
    translated_chunks = []
    for chunk in chunks:
        if chunk.strip():
            try:
                if len(chunk) > 4000:
                    sub_chunks = [chunk[i:i+4000] for i in range(0, len(chunk), 4000)]
                    t_sub = [translator.translate(sc) for sc in sub_chunks]
                    translated_chunks.append(" ".join(t_sub))
                else:
                    translated_chunks.append(translator.translate(chunk))
            except Exception as e:
                logger.warning(f"Translation chunk error: {e}")
                translated_chunks.append(chunk)
        else:
            translated_chunks.append("")
    return "\n".join(translated_chunks)


def ocr_extract_text_sync(image_or_pdf_path: str) -> str:
    """Skaner qilingan rasm yoki PDF dan OCR orqali matn tanish"""
    if not pytesseract:
        return "⚠️ OCR moduli o'rnatilmagan yoki Tesseract OCR binary topilmadi."

    ext = Path(image_or_pdf_path).suffix.lower()
    text_results = []
    
    if ext == ".pdf":
        doc = fitz.open(image_or_pdf_path)
        try:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=150)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                try:
                    page_text = pytesseract.image_to_string(img, lang="eng+rus")
                except Exception:
                    page_text = pytesseract.image_to_string(img)
                text_results.append(f"--- Sahifa {i+1} (OCR) ---\n{page_text}")
        finally:
            doc.close()
    else:
        try:
            img = Image.open(image_or_pdf_path)
            try:
                ocr_text = pytesseract.image_to_string(img, lang="eng+rus")
            except Exception:
                ocr_text = pytesseract.image_to_string(img)
            text_results.append(ocr_text)
        except Exception as e:
            logger.error(f"OCR Error: {e}")

    return "\n\n".join(text_results)


def pdf_to_images_sync(pdf_path: str, output_dir: str) -> List[str]:
    """PDF sahifalarini JPG rasmlarga ajratish"""
    doc = fitz.open(pdf_path)
    image_paths = []
    for i in range(len(doc)):
        page = doc[i]
        pix = page.get_pixmap(dpi=150)
        img_path = os.path.join(output_dir, f"page_{i+1}.jpg")
        pix.save(img_path)
        image_paths.append(img_path)
    doc.close()
    return image_paths


def pdf_extract_text_sync(pdf_path: str) -> str:
    """PDF dan matn ajratish"""
    doc = fitz.open(pdf_path)
    text = []
    for i, page in enumerate(doc):
        text.append(f"--- Sahifa {i+1} ---\n")
        text.append(page.get_text())
    doc.close()
    return "\n".join(text)


def pdf_compress_sync(input_pdf_path: str, output_pdf_path: str) -> str:
    """PDF hajmini siqish"""
    doc = fitz.open(input_pdf_path)
    doc.save(output_pdf_path, garbage=4, deflate=True, clean=True)
    doc.close()
    return output_pdf_path


# ==========================================
# TELEGRAM BOT HANDLERS & CONVERSATIONS
# ==========================================

ALBUM_BUFFERS: Dict[str, Dict] = {}
PENDING_IMAGE_TASKS: Dict[str, List[str]] = {}

def get_main_menu_keyboard():
    keyboard = [
        [
            InlineKeyboardButton("📸 Rasmlar / Albom -> PDF", callback_data="menu:img2pdf"),
            InlineKeyboardButton("📄 Hujjatlar -> PDF", callback_data="menu:doc2pdf"),
        ],
        [
            InlineKeyboardButton("📝 PDF -> Word (.docx)", callback_data="menu:pdf2docx"),
            InlineKeyboardButton("🖼 PDF -> Rasmlar (JPG)", callback_data="menu:pdf2img"),
        ],
        [
            InlineKeyboardButton("🔒 PDF Parol Boshqaruvi", callback_data="menu:pdfpass"),
            InlineKeyboardButton("✂️ Sahifalarni O'chirish", callback_data="menu:pdfdel"),
        ],
        [
            InlineKeyboardButton("💧 Suv Belgisi (Watermark)", callback_data="menu:watermark"),
            InlineKeyboardButton("✍️ Imzo / Muhr Qo'yish", callback_data="menu:stamp"),
        ],
        [
            InlineKeyboardButton("🗜 PDF Hajmini Siqish", callback_data="menu:compress"),
            InlineKeyboardButton("🔍 OCR Matn Tanish", callback_data="menu:ocr"),
        ],
        [
            InlineKeyboardButton("🌐 O'zbekcha Tarjima", callback_data="menu:translate"),
            InlineKeyboardButton("📖 Yordam & Yo'riqnoma", callback_data="menu:help"),
        ]
    ]
    return InlineKeyboardMarkup(keyboard)


async def start_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    welcome = (
        f"👋 Salom, <b>{user.first_name}</b>!\n\n"
        "🤖 <b>Super PDF Converter Bot</b>ga xush kelibsiz.\n\n"
        "👇 <i>Kerakli bo'limni tanlang yoki fayl/rasm yuboring:</i>"
    )
    await update.message.reply_text(welcome, parse_mode="HTML", reply_markup=get_main_menu_keyboard())


async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    help_text = (
        "📖 <b>Yordam yo'riqnomasi</b>\n\n"
        "1. <b>Albom rasmlar:</b> Bir nechta rasmni tanlab yuboring, bot ularni avtomatik bitta PDF qiladi.\n"
        "2. <b>Hujjatlar:</b> .docx, .xlsx, .pptx, .txt fayl yuboring — PDF bo'lib qaytadi.\n"
        "3. <b>PDF Fayllar:</b> PDF yuborsangiz, interaktiv menyu ochiladi. Parol o'rnatish, Word ga o'girish, Suv belgisi yoki Imzo qo'shishingiz mumkin.\n\n"
        "<i>Buyruqlar:</i>\n"
        "/start - Qayta ishga tushirish\n"
        "/help - Yo'riqroma\n"
        "/cancel - Harakatni bekor qilish"
    )
    await update.message.reply_text(help_text, parse_mode="HTML")


# ------------------------------------------
# ALBUM & PHOTO HANDLERS
# ------------------------------------------

async def prompt_image_fit_mode(bot, chat_id: int, photos: List[str], status_msg=None):
    """Foydalanuvchiga A4 o'lchamga moslash bo'yicha Inline Keyboard prompt chiqarish"""
    task_id = f"img_{datetime.now().strftime('%Y%m%d_%H%M%S_%f')}"
    PENDING_IMAGE_TASKS[task_id] = photos

    keyboard = [
        [
            InlineKeyboardButton("✅ Ha (A4 qog'ozga moslash)", callback_data=f"imgfit:a4:{task_id}"),
            InlineKeyboardButton("❌ Yo'q (Asl o'lchamda saqlash)", callback_data=f"imgfit:orig:{task_id}"),
        ]
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)

    text = (
        "📐 <b>Rasmlar PDF joylashuv rejimi:</b>\n\n"
        f"Menga <b>{len(photos)} ta rasm</b> yuborildi.\n"
        "Ushbu rasmlar A4 qog'oz o'lchamiga moslab (A4 Fit) joylashtirilsinmi?"
    )

    if status_msg:
        try:
            await status_msg.edit_text(text, parse_mode="HTML", reply_markup=reply_markup)
            return
        except Exception:
            pass

    await bot.send_message(chat_id=chat_id, text=text, parse_mode="HTML", reply_markup=reply_markup)


async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    msg = update.message
    
    # Agar foydalanuvchi "Imzo/Muhr" rejimida bo'lsa
    if context.user_data.get("awaiting_stamp_img"):
        file_id = context.user_data.pop("stamp_pdf_file_id", None)
        context.user_data["awaiting_stamp_img"] = False
        if file_id:
            await process_stamp_action(update, context, file_id, msg.photo[-1].file_id)
            return

    media_group_id = msg.media_group_id
    chat_id = update.effective_chat.id
    photo = msg.photo[-1]

    if media_group_id:
        if media_group_id not in ALBUM_BUFFERS:
            ALBUM_BUFFERS[media_group_id] = {
                "chat_id": chat_id,
                "photos": [],
                "timer_task": None,
                "status_msg": None,
            }
        buffer = ALBUM_BUFFERS[media_group_id]
        
        file = await photo.get_file()
        tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
        tmp_img.close()
        await file.download_to_drive(tmp_img.name)
        buffer["photos"].append(tmp_img.name)

        if buffer["status_msg"] is None:
            buffer["status_msg"] = await msg.reply_text("📥 Albom rasmlari qabul qilinmoqda...")

        if buffer["timer_task"]:
            buffer["timer_task"].cancel()

        buffer["timer_task"] = asyncio.create_task(process_album_after_delay(media_group_id, context))
        return

    # Yakka rasm
    file = await photo.get_file()
    status_msg = await msg.reply_text("📥 Rasm yuklanmoqda...")
    tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
    tmp_img.close()
    await file.download_to_drive(tmp_img.name)

    await prompt_image_fit_mode(context.bot, chat_id, [tmp_img.name], status_msg)


async def process_album_after_delay(media_group_id: str, context: ContextTypes.DEFAULT_TYPE):
    await asyncio.sleep(1.8)
    if media_group_id not in ALBUM_BUFFERS:
        return

    buffer = ALBUM_BUFFERS.pop(media_group_id)
    chat_id = buffer["chat_id"]
    photos = buffer["photos"]
    status_msg = buffer["status_msg"]

    await prompt_image_fit_mode(context.bot, chat_id, photos, status_msg)


# ------------------------------------------
# DOCUMENT HANDLERS (.docx, .xlsx, .pptx, .txt, .pdf)
# ------------------------------------------

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    msg = update.message
    doc = msg.document
    if not doc:
        return

    filename = doc.file_name or "document"
    ext = Path(filename).suffix.lower()

    # Image document
    if is_image_file(filename):
        status_msg = await msg.reply_text(f"📥 Rasm yuklanmoqda... {render_progress_bar(30)}")
        file = await doc.get_file()
        with tempfile.TemporaryDirectory() as tmp_dir:
            input_path = os.path.join(tmp_dir, filename)
            output_pdf = os.path.join(tmp_dir, f"{Path(filename).stem}.pdf")
            await file.download_to_drive(input_path)
            await status_msg.edit_text(f"⚡ PDF yaratilmoqda... {render_progress_bar(80)}")
            try:
                await asyncio.to_thread(images_to_pdf_sync, [input_path], output_pdf)
                with open(output_pdf, "rb") as f:
                    await msg.reply_document(document=f, filename=f"{Path(filename).stem}.pdf", caption="✅ PDF tayyor!")
                await status_msg.delete()
            except Exception as e:
                await status_msg.edit_text(f"❌ Xatolik: {e}")
        return

    # Word (.docx)
    if is_docx_file(filename):
        status_msg = await msg.reply_text(f"📥 Word fayli yuklanmoqda... {render_progress_bar(25)}")
        file = await doc.get_file()
        with tempfile.TemporaryDirectory() as tmp_dir:
            input_path = os.path.join(tmp_dir, filename)
            output_pdf = os.path.join(tmp_dir, f"{Path(filename).stem}.pdf")
            await file.download_to_drive(input_path)
            await status_msg.edit_text(f"⚡ PDF shakllantirilmoqda... {render_progress_bar(75)}")
            try:
                await asyncio.to_thread(docx_to_pdf_sync, input_path, output_pdf)
                with open(output_pdf, "rb") as f:
                    await msg.reply_document(document=f, filename=f"{Path(filename).stem}.pdf", caption="✅ Word -> PDF tayyor!")
                await status_msg.delete()
            except Exception as e:
                await status_msg.edit_text(f"❌ Word o'girishda xatolik: {e}")
        return

    # Excel (.xlsx)
    if is_excel_file(filename):
        status_msg = await msg.reply_text(f"📥 Excel fayli yuklanmoqda... {render_progress_bar(25)}")
        file = await doc.get_file()
        with tempfile.TemporaryDirectory() as tmp_dir:
            input_path = os.path.join(tmp_dir, filename)
            output_pdf = os.path.join(tmp_dir, f"{Path(filename).stem}.pdf")
            await file.download_to_drive(input_path)
            await status_msg.edit_text(f"⚡ Excel jadvallari PDF qilinmoqda... {render_progress_bar(75)}")
            try:
                await asyncio.to_thread(excel_to_pdf_sync, input_path, output_pdf)
                with open(output_pdf, "rb") as f:
                    await msg.reply_document(document=f, filename=f"{Path(filename).stem}.pdf", caption="✅ Excel -> PDF tayyor!")
                await status_msg.delete()
            except Exception as e:
                await status_msg.edit_text(f"❌ Excel o'girishda xatolik: {e}")
        return

    # PowerPoint (.pptx)
    if is_pptx_file(filename):
        status_msg = await msg.reply_text(f"📥 Prezentatsiya yuklanmoqda... {render_progress_bar(25)}")
        file = await doc.get_file()
        with tempfile.TemporaryDirectory() as tmp_dir:
            input_path = os.path.join(tmp_dir, filename)
            output_pdf = os.path.join(tmp_dir, f"{Path(filename).stem}.pdf")
            await file.download_to_drive(input_path)
            await status_msg.edit_text(f"⚡ Slaydlar PDF qilinmoqda... {render_progress_bar(75)}")
            try:
                await asyncio.to_thread(pptx_to_pdf_sync, input_path, output_pdf)
                with open(output_pdf, "rb") as f:
                    await msg.reply_document(document=f, filename=f"{Path(filename).stem}.pdf", caption="✅ PPTX -> PDF tayyor!")
                await status_msg.delete()
            except Exception as e:
                await status_msg.edit_text(f"❌ Prezentatsiya o'girishda xatolik: {e}")
        return

    # Text (.txt, .md, .log, .json ...)
    if is_text_file(filename):
        status_msg = await msg.reply_text(f"📥 Matn fayli yuklanmoqda... {render_progress_bar(30)}")
        file = await doc.get_file()
        with tempfile.TemporaryDirectory() as tmp_dir:
            input_path = os.path.join(tmp_dir, filename)
            output_pdf = os.path.join(tmp_dir, f"{Path(filename).stem}.pdf")
            await file.download_to_drive(input_path)
            await status_msg.edit_text(f"⚡ PDF formatiga o'tkazilmoqda... {render_progress_bar(85)}")
            try:
                with open(input_path, "rb") as f:
                    raw_bytes = f.read()
                detected = chardet.detect(raw_bytes)
                encoding = detected.get("encoding") or "utf-8"
                try:
                    text_content = raw_bytes.decode(encoding)
                except Exception:
                    text_content = raw_bytes.decode("utf-8", errors="ignore")
                
                await asyncio.to_thread(text_to_pdf_sync, text_content, output_pdf, Path(filename).stem)
                with open(output_pdf, "rb") as f:
                    await msg.reply_document(document=f, filename=f"{Path(filename).stem}.pdf", caption="✅ Matn -> PDF tayyor!")
                await status_msg.delete()
            except Exception as e:
                await status_msg.edit_text(f"❌ Matn o'girishda xatolik: {e}")
        return

    # PDF File (Toolkit Menu)
    if ext == ".pdf":
        file_id = doc.file_id
        keyboard = [
            [
                InlineKeyboardButton("📝 PDF -> Word (.docx)", callback_data=f"pdf:to_docx:{file_id}"),
                InlineKeyboardButton("🖼 PDF -> Rasmlar", callback_data=f"pdf:to_img:{file_id}"),
            ],
            [
                InlineKeyboardButton("🌐 O'zbekchaga tarjima qilish", callback_data=f"pdf:translate:{file_id}"),
                InlineKeyboardButton("🔍 OCR (Skaner matnini olish)", callback_data=f"pdf:ocr:{file_id}"),
            ],
            [
                InlineKeyboardButton("🔒 Parol qo'yish", callback_data=f"pdf:encrypt:{file_id}"),
                InlineKeyboardButton("🔓 Parolni yechish", callback_data=f"pdf:decrypt:{file_id}"),
            ],
            [
                InlineKeyboardButton("✂️ Sahifalarni o'chirish", callback_data=f"pdf:del_pages:{file_id}"),
                InlineKeyboardButton("💧 Suv belgisi", callback_data=f"pdf:watermark:{file_id}"),
            ],
            [
                InlineKeyboardButton("✍️ Imzo / Muhr qo'yish", callback_data=f"pdf:stamp:{file_id}"),
                InlineKeyboardButton("🗜 Hajmini siqish", callback_data=f"pdf:compress:{file_id}"),
            ],
            [
                InlineKeyboardButton("📄 Matnni ajratish", callback_data=f"pdf:to_txt:{file_id}"),
            ]
        ]
        reply_markup = InlineKeyboardMarkup(keyboard)
        await msg.reply_text(
            f"📄 <b>{filename}</b> qabul qilindi.\n\n"
            "Ushbu PDF bo'yicha amalnigizni tanlang:",
            parse_mode="HTML",
            reply_markup=reply_markup,
        )
        return

    await msg.reply_text(
        f"⚠️ <b>{ext}</b> formati qo'llab-quvvatlanmaydi."
    )


# ------------------------------------------
# PDF TOOLKIT CALLBACK HANDLER
# ------------------------------------------

async def handle_pdf_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()

    data = query.data
    msg = query.message

    if data.startswith("menu:"):
        action = data.split(":", 1)[1]
        back_kb = InlineKeyboardMarkup([[InlineKeyboardButton("🏠 Bosh Menyuga Qaytish", callback_data="menu:home")]])
        
        if action == "home":
            welcome = (
                "👋 Salom!\n\n"
                "🤖 <b>Super PDF Converter Bot</b>ga xush kelibsiz.\n\n"
                "👇 <i>Kerakli bo'limni tanlang yoki fayl/rasm yuboring:</i>"
            )
            await msg.edit_text(welcome, parse_mode="HTML", reply_markup=get_main_menu_keyboard())
            return

        menu_texts = {
            "img2pdf": "📸 <b>Rasmlar / Albomlar -> PDF</b>\n\nMenga 1 ta yoki bir nechta rasm yuboring! Bot ularni avtomatik yagona ko'p sahifali PDF formatiga o'giradi.",
            "doc2pdf": "📄 <b>Hujjatlar -> PDF</b>\n\nMenga Word (.docx), Excel (.xlsx), PowerPoint (.pptx) yoki Matn (.txt) faylini yuboring — bot uni PDF ga aylantiradi.",
            "pdf2docx": "📝 <b>PDF -> Word (.docx)</b>\n\nMenga PDF faylingizni yuboring va menyudan <b>'PDF -> Word (.docx)'</b> tugmasini bosing.",
            "pdf2img": "🖼 <b>PDF -> Rasmlar (JPG)</b>\n\nMenga PDF faylingizni yuboring va menyudan <b>'PDF -> Rasmlar'</b> tugmasini bosing.",
            "pdfpass": "🔒 <b>PDF Parol Boshqaruvi</b>\n\nMenga PDF faylingizni yuboring va menyudan <b>'Parol qo'yish'</b> yoki <b>'Parolni yechish'</b> tugmasini bosing.",
            "pdfdel": "✂️ <b>Sahifalarni O'chirish</b>\n\nMenga PDF faylingizni yuboring va menyudan <b>'Sahifalarni o'chirish'</b> tugmasini bosing.",
            "watermark": "💧 <b>Suv Belgisi (Watermark)</b>\n\nMenga PDF faylingizni yuboring va menyudan <b>'Suv belgisi'</b> tugmasini bosing.",
            "stamp": "✍️ <b>Imzo / Muhr Qo'yish</b>\n\nMenga PDF faylingizni yuboring va menyudan <b>'Imzo / Muhr qo'yish'</b> tugmasini bosing.",
            "compress": "🗜 <b>PDF Hajmini Siqish</b>\n\nMenga PDF faylingizni yuboring va menyudan <b>'Hajmini siqish'</b> tugmasini bosing.",
            "ocr": "🔍 <b>OCR Matn Tanish</b>\n\nSkaner qilingan rasm yoki PDF yuboring va menyudan <b>'OCR (Skaner matnini olish)'</b> tugmasini bosing.",
            "translate": "🌐 <b>O'zbekcha Tarjima</b>\n\nPDF yuboring va menyudan <b>'O'zbekchaga tarjima qilish'</b> tugmasini bosing.",
            "help": "📖 <b>Yordam yo'riqnomasi</b>\n\n1. Har qanday rasm, Word, Excel, PPTX, TXT yoki PDF yuboring.\n2. PDF yuborsangiz, barcha uskunalar menyusi avtomatik chiqadi."
        }
        text = menu_texts.get(action, "👇 Fayl yoki rasm yuborishingiz mumkin!")
        await msg.edit_text(text, parse_mode="HTML", reply_markup=back_kb)
        return

    if data.startswith("imgfit:"):
        parts = data.split(":", 2)
        fit_mode = parts[1]
        task_id = parts[2]
        
        photos = PENDING_IMAGE_TASKS.pop(task_id, None)
        if not photos:
            await msg.edit_text("⚠️ Rasmlar topilmadi yoki sessiya muddati tugadi.")
            return

        await msg.edit_text(f"⚡ PDF yaratilmoqda... {render_progress_bar(60)}")
        with tempfile.TemporaryDirectory() as tmp_dir:
            output_pdf = os.path.join(tmp_dir, "converted.pdf")
            try:
                if fit_mode == "a4":
                    await asyncio.to_thread(images_to_pdf_a4_sync, photos, output_pdf)
                else:
                    await asyncio.to_thread(images_to_pdf_sync, photos, output_pdf)
                
                await msg.edit_text(f"✅ Tayyor! {render_progress_bar(100)}")
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                with open(output_pdf, "rb") as f:
                    await context.bot.send_document(
                        chat_id=msg.chat_id,
                        document=f,
                        filename=f"converted_{timestamp}.pdf",
                        caption=f"✅ {len(photos)} ta rasmdan iborat PDF tayyor bo'ldi!",
                    )
                await msg.delete()
            except Exception as e:
                logger.error(f"Image PDF fit error: {e}")
                await msg.edit_text(f"❌ PDF yaratishda xatolik: {e}")
            finally:
                for p in photos:
                    if os.path.exists(p):
                        try:
                            os.remove(p)
                        except Exception:
                            pass
        return

    if not data.startswith("pdf:"):
        return

    parts = data.split(":", 2)
    action = parts[1]
    file_id = parts[2]

    # 1. PDF -> Word (.docx)
    if action == "to_docx":
        await msg.edit_text(f"⏳ PDF Word (.docx) ga o'girilmoqda... {render_progress_bar(40)}")
        try:
            telegram_file = await context.bot.get_file(file_id)
            with tempfile.TemporaryDirectory() as tmp_dir:
                pdf_path = os.path.join(tmp_dir, "input.pdf")
                docx_path = os.path.join(tmp_dir, "converted.docx")
                await telegram_file.download_to_drive(pdf_path)
                
                await msg.edit_text(f"⚡ Word struktura tuzilmoqda... {render_progress_bar(80)}")
                await asyncio.to_thread(pdf_to_docx_sync, pdf_path, docx_path)
                
                with open(docx_path, "rb") as f:
                    await context.bot.send_document(
                        chat_id=msg.chat_id,
                        document=f,
                        filename="converted.docx",
                        caption="✅ PDF Word ga muvaffaqiyatli o'girildi!",
                    )
                await msg.delete()
        except Exception as e:
            logger.error(f"PDF to docx error: {e}")
            await msg.edit_text(f"❌ Word o'girishda xatolik: {e}")

    # 2. PDF -> O'zbekcha Tarjima
    elif action == "translate":
        await msg.edit_text(f"🌐 PDF matni o'zbek tiliga tarjima qilinmoqda... {render_progress_bar(40)}")
        try:
            telegram_file = await context.bot.get_file(file_id)
            with tempfile.TemporaryDirectory() as tmp_dir:
                pdf_path = os.path.join(tmp_dir, "input.pdf")
                output_pdf = os.path.join(tmp_dir, "translated.pdf")
                await telegram_file.download_to_drive(pdf_path)
                
                raw_text = await asyncio.to_thread(pdf_extract_text_sync, pdf_path)
                await msg.edit_text(f"⚡ Matn o'zbek tiliga o'girilmoqda... {render_progress_bar(70)}")
                translated_text = await asyncio.to_thread(translate_text_sync, raw_text, "uz")
                
                await asyncio.to_thread(text_to_pdf_sync, translated_text, output_pdf, "Tarjima qilingan hujjat")
                
                with open(output_pdf, "rb") as f:
                    await context.bot.send_document(
                        chat_id=msg.chat_id,
                        document=f,
                        filename="translated_uz.pdf",
                        caption="✅ PDF matni o'zbek tiliga tarjima qilindi!",
                    )
                await msg.delete()
        except Exception as e:
            logger.error(f"PDF translate error: {e}")
            await msg.edit_text(f"❌ Tarjimada xatolik: {e}")

    # 3. PDF -> OCR Matn Tanish
    elif action == "ocr":
        await msg.edit_text(f"🔍 Skaner PDF OCR yordamida o'qilmoqda... {render_progress_bar(50)}")
        try:
            telegram_file = await context.bot.get_file(file_id)
            with tempfile.TemporaryDirectory() as tmp_dir:
                pdf_path = os.path.join(tmp_dir, "input.pdf")
                await telegram_file.download_to_drive(pdf_path)
                
                ocr_text = await asyncio.to_thread(ocr_extract_text_sync, pdf_path)
                if not ocr_text.strip():
                    await msg.edit_text("⚠️ OCR yordamida matn topilmadi.")
                    return
                
                txt_path = os.path.join(tmp_dir, "ocr_text.txt")
                with open(txt_path, "w", encoding="utf-8") as f:
                    f.write(ocr_text)
                
                with open(txt_path, "rb") as f:
                    await context.bot.send_document(
                        chat_id=msg.chat_id,
                        document=f,
                        filename="ocr_extracted.txt",
                        caption="🔍 OCR orqali aniqlangan matn!",
                    )
                await msg.delete()
        except Exception as e:
            logger.error(f"PDF OCR error: {e}")
            await msg.edit_text(f"❌ OCR amalda xatolik: {e}")

    # 2. PDF -> Rasmlar (JPG)
    elif action == "to_img":
        await msg.edit_text(f"⏳ Sahifalar rasmlarga ajratilmoqda... {render_progress_bar(50)}")
        try:
            telegram_file = await context.bot.get_file(file_id)
            with tempfile.TemporaryDirectory() as tmp_dir:
                pdf_path = os.path.join(tmp_dir, "input.pdf")
                await telegram_file.download_to_drive(pdf_path)
                
                img_paths = await asyncio.to_thread(pdf_to_images_sync, pdf_path, tmp_dir)
                if not img_paths:
                    await msg.edit_text("❌ PDF ichida sahifa topilmadi.")
                    return
                
                await msg.edit_text(f"📤 {len(img_paths)} ta rasm yuborilmoqda... {render_progress_bar(90)}")
                for i in range(0, min(len(img_paths), 30), 10):
                    chunk = img_paths[i:i+10]
                    media_group = [InputMediaDocument(media=open(p, "rb")) for p in chunk]
                    await context.bot.send_media_group(chat_id=msg.chat_id, media=media_group)
                await msg.delete()
        except Exception as e:
            await msg.edit_text(f"❌ Rasmlarga ajratishda xatolik: {e}")

    # 3. PDF Encrypt (Parol qo'yish)
    elif action == "encrypt":
        context.user_data["action_pdf_file_id"] = file_id
        context.user_data["pending_action"] = "encrypt"
        await msg.edit_text(
            "🔒 <b>PDF ga parol o'rnatish</b>\n\n"
            "Iltimos, PDF fayl uchun o'rnatmoqchi bo'lgan parolingizni yozib yuboring:",
            parse_mode="HTML",
        )

    # 4. PDF Decrypt (Parol olib tashlash)
    elif action == "decrypt":
        context.user_data["action_pdf_file_id"] = file_id
        context.user_data["pending_action"] = "decrypt"
        await msg.edit_text(
            "🔓 <b>PDF parolini yechish</b>\n\n"
            "Iltimos, ushbu PDF faylning parolini yozib yuboring:",
            parse_mode="HTML",
        )

    # 5. Delete Pages (Sahifalarni o'chirish)
    elif action == "del_pages":
        context.user_data["action_pdf_file_id"] = file_id
        context.user_data["pending_action"] = "del_pages"
        await msg.edit_text(
            "✂️ <b>Sahifalarni o'chirish</b>\n\n"
            "O'chirmoqchi bo megan sahifa raqamlarini vergul bilan yuboring:\n"
            "<i>Masalan: 2, 5, 8</i>",
            parse_mode="HTML",
        )

    # 6. Watermark (Suv belgisi)
    elif action == "watermark":
        context.user_data["action_pdf_file_id"] = file_id
        context.user_data["pending_action"] = "watermark"
        await msg.edit_text(
            "💧 <b>Suv belgisi (Watermark) qo'shish</b>\n\n"
            "PDF sahifalariga yoziladigan matnni yuboring:\n"
            "<i>Masalan: MAXFIY yoki Shaxsiy</i>",
            parse_mode="HTML",
        )

    # 7. Stamp / Imzo
    elif action == "stamp":
        context.user_data["stamp_pdf_file_id"] = file_id
        context.user_data["awaiting_stamp_img"] = True
        await msg.edit_text(
            "✍️ <b>Imzo yoki Muhr joylashtirish</b>\n\n"
            "Iltimos, PDF-ga qo'ymoqchi bo'lgan <b>shaffof PNG imzo/muhr rasmini</b> yuboring:",
            parse_mode="HTML",
        )

    # 8. Compress
    elif action == "compress":
        await msg.edit_text(f"⚡ PDF siqilmoqda... {render_progress_bar(45)}")
        try:
            telegram_file = await context.bot.get_file(file_id)
            with tempfile.TemporaryDirectory() as tmp_dir:
                pdf_path = os.path.join(tmp_dir, "input.pdf")
                out_path = os.path.join(tmp_dir, "compressed.pdf")
                await telegram_file.download_to_drive(pdf_path)
                orig_size = os.path.getsize(pdf_path)
                
                await asyncio.to_thread(pdf_compress_sync, pdf_path, out_path)
                comp_size = os.path.getsize(out_path)
                
                orig_mb = round(orig_size / (1024 * 1024), 2)
                comp_mb = round(comp_size / (1024 * 1024), 2)
                
                with open(out_path, "rb") as f:
                    await context.bot.send_document(
                        chat_id=msg.chat_id,
                        document=f,
                        filename="compressed.pdf",
                        caption=f"✅ PDF hajmi siqildi!\n\n📊 Avvalgi: {orig_mb} MB\n📉 Yangi: {comp_mb} MB",
                    )
                await msg.delete()
        except Exception as e:
            await msg.edit_text(f"❌ PDF siqishda xatolik: {e}")

    # 9. Extract Text
    elif action == "to_txt":
        await msg.edit_text(f"⏳ Matn ajratib olinmoqda... {render_progress_bar(50)}")
        try:
            telegram_file = await context.bot.get_file(file_id)
            with tempfile.TemporaryDirectory() as tmp_dir:
                pdf_path = os.path.join(tmp_dir, "input.pdf")
                await telegram_file.download_to_drive(pdf_path)
                
                text = await asyncio.to_thread(pdf_extract_text_sync, pdf_path)
                if not text.strip():
                    await msg.edit_text("⚠️ PDF ichida matn topilmadi.")
                    return
                
                txt_path = os.path.join(tmp_dir, "extracted.txt")
                with open(txt_path, "w", encoding="utf-8") as f:
                    f.write(text)
                
                with open(txt_path, "rb") as f:
                    await context.bot.send_document(
                        chat_id=msg.chat_id,
                        document=f,
                        filename="extracted_text.txt",
                        caption="✅ Ajratib olingan matn!",
                    )
                await msg.delete()
        except Exception as e:
            await msg.edit_text(f"❌ Matn ajratishda xatolik: {e}")


# ------------------------------------------
# USER TEXT RESPONSE HANDLER (For Pending Actions)
# ------------------------------------------

async def handle_user_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_text = update.message.text.strip()
    pending_action = context.user_data.pop("pending_action", None)
    file_id = context.user_data.pop("action_pdf_file_id", None)

    if not pending_action or not file_id:
        await update.message.reply_text(
            "Iltimos, PDF ga aylantirish uchun rasm yoki fayl yuboring.\n"
            "Yordam uchun /help buyrug'ini ishlating."
        )
        return

    status_msg = await update.message.reply_text(f"⚡ Bajarilmoqda... {render_progress_bar(40)}")

    try:
        telegram_file = await context.bot.get_file(file_id)
        with tempfile.TemporaryDirectory() as tmp_dir:
            input_pdf = os.path.join(tmp_dir, "input.pdf")
            output_pdf = os.path.join(tmp_dir, "processed.pdf")
            await telegram_file.download_to_drive(input_pdf)

            if pending_action == "encrypt":
                await asyncio.to_thread(pdf_encrypt_sync, input_pdf, output_pdf, user_text)
                caption = "🔒 PDF parollangan holda saqlandi!"
                filename = "protected.pdf"

            elif pending_action == "decrypt":
                await asyncio.to_thread(pdf_decrypt_sync, input_pdf, output_pdf, user_text)
                caption = "🔓 PDF paroli olib tashlandi!"
                filename = "unlocked.pdf"

            elif pending_action == "del_pages":
                pages = [int(p.strip()) for p in re.split(r"[,; ]+", user_text) if p.strip().isdigit()]
                if not pages:
                    await status_msg.edit_text("❌ Yaroqli sahifa raqamlari kiritilmadi.")
                    return
                await asyncio.to_thread(pdf_delete_pages_sync, input_pdf, output_pdf, pages)
                caption = f"✂️ {len(pages)} ta sahifa o'chirildi!"
                filename = "modified.pdf"

            elif pending_action == "watermark":
                await asyncio.to_thread(pdf_add_watermark_sync, input_pdf, output_pdf, user_text)
                caption = "💧 Suv belgisi muvaffaqiyatli qo'shildi!"
                filename = "watermarked.pdf"

            else:
                await status_msg.edit_text("❌ Noma'lum amal.")
                return

            await status_msg.edit_text(f"✅ Tayyor! {render_progress_bar(100)}")
            with open(output_pdf, "rb") as f:
                await update.message.reply_document(
                    document=f,
                    filename=filename,
                    caption=caption,
                )
            await status_msg.delete()

    except Exception as e:
        logger.error(f"Pending action error ({pending_action}): {e}")
        await status_msg.edit_text(f"❌ Xatolik yuz berdi: {e}")


async def process_stamp_action(update: Update, context: ContextTypes.DEFAULT_TYPE, pdf_file_id: str, photo_file_id: str):
    status_msg = await update.message.reply_text(f"⚡ Imzo PDF-ga joylashtirilmoqda... {render_progress_bar(50)}")
    try:
        pdf_file = await context.bot.get_file(pdf_file_id)
        photo_file = await context.bot.get_file(photo_file_id)
        
        with tempfile.TemporaryDirectory() as tmp_dir:
            input_pdf = os.path.join(tmp_dir, "input.pdf")
            stamp_img = os.path.join(tmp_dir, "stamp.png")
            output_pdf = os.path.join(tmp_dir, "stamped.pdf")
            
            await pdf_file.download_to_drive(input_pdf)
            await photo_file.download_to_drive(stamp_img)
            
            await asyncio.to_thread(pdf_add_stamp_sync, input_pdf, output_pdf, stamp_img, 1)
            
            await status_msg.edit_text(f"✅ Tayyor! {render_progress_bar(100)}")
            with open(output_pdf, "rb") as f:
                await update.message.reply_document(
                    document=f,
                    filename="stamped_document.pdf",
                    caption="✍️ Imzo/Muhr PDF-ga joylashtirildi!",
                )
            await status_msg.delete()
    except Exception as e:
        logger.error(f"Stamp process error: {e}")
        await status_msg.edit_text(f"❌ Imzo qo'shishda xatolik: {e}")


from http.server import HTTPServer, BaseHTTPRequestHandler
import threading

class HealthCheckHandler(BaseHTTPRequestHandler):
    def do_GET(self):
        self.send_response(200)
        self.send_header("Content-type", "text/plain")
        self.end_headers()
        self.wfile.write(b"OK")
    def log_message(self, format, *args):
        pass

def start_health_check_server():
    port_str = os.getenv("PORT")
    if port_str:
        try:
            port = int(port_str)
            server = HTTPServer(("0.0.0.0", port), HealthCheckHandler)
            thread = threading.Thread(target=server.serve_forever, daemon=True)
            thread.start()
            logger.info(f"Render Health Check HTTP server {port}-portda ishga tushdi.")
        except Exception as e:
            logger.warning(f"Health Check Server xatosi: {e}")


def main():
    if TOKEN == "YOUR_BOT_TOKEN_HERE" or not TOKEN:
        print("=" * 50)
        print("XATOLIK: BOT_TOKEN o'rnatilmagan!")
        print("=" * 50)
        return

    # Render port binding uchun HTTP Health Check serverini ishga tushirish
    start_health_check_server()

    # Python 3.10+ / 3.14 Event loop tayyorlash
    try:
        loop = asyncio.get_event_loop()
    except RuntimeError:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)

    app = Application.builder().token(TOKEN).build()

    app.add_handler(CommandHandler("start", start_command))
    app.add_handler(CommandHandler("help", help_command))

    app.add_handler(MessageHandler(filters.PHOTO & ~filters.COMMAND, handle_photo))
    app.add_handler(MessageHandler(filters.Document.ALL & ~filters.COMMAND, handle_document))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_user_text))
    app.add_handler(CallbackQueryHandler(handle_pdf_callback))

    print("🚀 Super PDF Converter Bot ishga tushdi...")
    print("To'xtatish uchun Ctrl+C bosing.")
    app.run_polling(allowed_updates=Update.ALL_TYPES)


if __name__ == "__main__":
    main()
